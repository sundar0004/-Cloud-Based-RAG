from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_CONNECTOR, MSO_SHAPE
from pptx.util import Inches, Pt

OUT_FILE = "docs/MSc_RAG_Project_Vels_University.pptx"
STUDENT_NAME = "Sundar A"
COLLEGE_ID = "24207105"
EMAIL = "sunadrsundar0004@gmail.com"
PHONE = "7904117804"
DEPARTMENT = "MSc Computer Science"
REFERENCE_NAME = "Akila"
GITHUB_REPO = "https://github.com/sundar0004/-Cloud-Based-RAG.git"

PALETTE = [
    RGBColor(8, 45, 108),
    RGBColor(0, 95, 115),
    RGBColor(18, 103, 130),
    RGBColor(37, 99, 235),
    RGBColor(20, 83, 45),
    RGBColor(146, 64, 14),
]


def style_slide(slide, title, idx):
    color = PALETTE[idx % len(PALETTE)]
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.65))
    band.fill.solid()
    band.fill.fore_color.rgb = color
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.4), Inches(0.08), Inches(12.5), Inches(0.45))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.text = title
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = RGBColor(255, 255, 255)


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(11, 44, 111)
    bg.line.fill.background()

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(5.8), Inches(13.33), Inches(1.7))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(37, 99, 235)
    accent.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.8), Inches(1.2), Inches(11.9), Inches(3.0)).text_frame
    title.text = "Design and Experimental Evaluation of a Cloud-Based RAG Framework"
    title.paragraphs[0].font.size = Pt(38)
    title.paragraphs[0].font.bold = True
    title.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)

    sub = slide.shapes.add_textbox(Inches(0.8), Inches(4.5), Inches(11.9), Inches(1.0)).text_frame
    sub.text = "Reducing Hallucination in Large Language Models"
    sub.paragraphs[0].font.size = Pt(24)
    sub.paragraphs[0].font.color.rgb = RGBColor(215, 228, 255)

    foot = slide.shapes.add_textbox(Inches(0.8), Inches(6.2), Inches(12.0), Inches(0.9)).text_frame
    foot.text = f"MSc Final Year Project | Vels University | {date.today().isoformat()}"
    foot.paragraphs[0].font.size = Pt(18)
    foot.paragraphs[0].font.bold = True
    foot.paragraphs[0].font.color.rgb = RGBColor(255, 255, 255)


def add_bullets_slide(prs, title, bullets, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide, title, idx)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.55), Inches(1.0), Inches(12.2), Inches(6.0))
    panel.fill.solid()
    panel.fill.fore_color.rgb = RGBColor(247, 250, 255)
    panel.line.color.rgb = RGBColor(198, 215, 242)

    box = slide.shapes.add_textbox(Inches(0.95), Inches(1.4), Inches(11.4), Inches(5.3))
    tf = box.text_frame
    tf.clear()
    for i, b in enumerate(bullets):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = f"• {b}"
        p.font.size = Pt(23)
        p.font.color.rgb = RGBColor(20, 30, 50)


def add_two_column_slide(prs, title, left_title, left_points, right_title, right_points, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide, title, idx)

    left = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.6), Inches(1.1), Inches(6.1), Inches(5.9))
    left.fill.solid(); left.fill.fore_color.rgb = RGBColor(239, 246, 255)
    left.line.color.rgb = RGBColor(147, 197, 253)

    right = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.7), Inches(1.1), Inches(6.1), Inches(5.9))
    right.fill.solid(); right.fill.fore_color.rgb = RGBColor(236, 253, 245)
    right.line.color.rgb = RGBColor(110, 231, 183)

    ltf = slide.shapes.add_textbox(Inches(0.9), Inches(1.35), Inches(5.5), Inches(5.3)).text_frame
    ltf.text = left_title
    ltf.paragraphs[0].font.size = Pt(24)
    ltf.paragraphs[0].font.bold = True
    for item in left_points:
        p = ltf.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(19)

    rtf = slide.shapes.add_textbox(Inches(7.0), Inches(1.35), Inches(5.5), Inches(5.3)).text_frame
    rtf.text = right_title
    rtf.paragraphs[0].font.size = Pt(24)
    rtf.paragraphs[0].font.bold = True
    for item in right_points:
        p = rtf.add_paragraph(); p.text = f"• {item}"; p.font.size = Pt(19)


def add_architecture_slide(prs, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide, "Complete Project Architecture", idx)

    def node(x, y, w, h, label, fill):
        shp = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
        shp.line.color.rgb = RGBColor(37, 99, 235)
        tf = shp.text_frame
        tf.text = label
        tf.paragraphs[0].font.size = Pt(13)
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.color.rgb = RGBColor(20, 30, 50)
        return shp

    ui = node(0.6, 1.4, 1.8, 0.9, "User / Web UI", RGBColor(224, 242, 254))
    api = node(2.8, 1.4, 2.0, 0.9, "FastAPI\nBackend", RGBColor(219, 234, 254))
    retr = node(5.2, 1.4, 2.0, 0.9, "Retriever", RGBColor(224, 231, 255))
    faiss = node(7.6, 1.4, 2.0, 0.9, "FAISS\nVector DB", RGBColor(226, 232, 240))
    llm = node(10.0, 1.4, 2.3, 0.9, "Local LLM\nGenerator", RGBColor(220, 252, 231))

    ingest = node(2.0, 4.2, 2.3, 0.9, "Document\nIngestion", RGBColor(254, 243, 199))
    embed = node(4.8, 4.2, 2.3, 0.9, "Embeddings", RGBColor(254, 240, 138))
    meta = node(7.6, 4.2, 2.3, 0.9, "Chunk Metadata", RGBColor(254, 226, 226))

    arrows = [
        (2.4, 1.85, 2.8, 1.85),
        (4.8, 1.85, 5.2, 1.85),
        (7.2, 1.85, 7.6, 1.85),
        (9.6, 1.85, 10.0, 1.85),
        (4.3, 4.65, 4.8, 4.65),
        (7.1, 4.65, 7.6, 4.65),
        (6.0, 4.2, 8.3, 2.3),
    ]
    for x1, y1, x2, y2 in arrows:
        conn = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
        conn.line.color.rgb = RGBColor(37, 99, 235)
        conn.line.width = Pt(2)

    note = slide.shapes.add_textbox(Inches(0.7), Inches(5.7), Inches(12.0), Inches(1.0)).text_frame
    note.text = "Flow: Documents -> Ingestion -> Embeddings -> FAISS -> Retrieval -> Context Injection -> Local LLM -> Answer"
    note.paragraphs[0].font.size = Pt(16)
    note.paragraphs[0].font.bold = True


def add_image_slide(prs, title, image_path, caption, idx):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    style_slide(slide, title, idx)

    panel = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.95), Inches(12.3), Inches(5.95))
    panel.fill.solid(); panel.fill.fore_color.rgb = RGBColor(248, 250, 252)
    panel.line.color.rgb = RGBColor(203, 213, 225)

    if Path(image_path).exists():
        slide.shapes.add_picture(image_path, Inches(0.8), Inches(1.2), Inches(11.7), Inches(5.3))

    cap = slide.shapes.add_textbox(Inches(0.8), Inches(6.7), Inches(11.8), Inches(0.5)).text_frame
    cap.text = caption
    cap.paragraphs[0].font.size = Pt(14)
    cap.paragraphs[0].font.color.rgb = RGBColor(51, 65, 85)


def build_deck():
    prs = Presentation()

    add_title_slide(prs)  # 1
    add_bullets_slide(prs, "Presenter Details", [
        f"Student Name: {STUDENT_NAME}",
        f"College ID: {COLLEGE_ID}",
        f"Department: {DEPARTMENT}",
        f"Reference Name: {REFERENCE_NAME}",
        f"Email: {EMAIL}",
        f"Phone: {PHONE}",
        f"GitHub Repository: {GITHUB_REPO}",
    ], 1)  # 2

    add_bullets_slide(prs, "Agenda", [
        "Problem and motivation",
        "Architecture and implementation",
        "Experiment strategy and metrics",
        "Timeline, risks, and references",
        "Frontend and backend demonstration",
    ], 2)  # 3

    add_bullets_slide(prs, "Problem Statement", [
        "LLMs can produce fluent but factually incorrect outputs",
        "Hallucination limits enterprise trust and adoption",
        "Need evidence-grounded generation from local documents",
        "RAG is used to improve reliability",
    ], 3)  # 4

    add_bullets_slide(prs, "Aim and Objectives", [
        "Design a cloud-ready RAG framework",
        "Compare baseline vs retrieval-augmented responses",
        "Evaluate factuality, hallucination, retrieval quality, and latency",
        "Build complete API + web UI + documentation package",
    ], 4)  # 5

    add_bullets_slide(prs, "Literature Foundation", [
        "Transformers: Attention Is All You Need (2017)",
        "Few-shot LLM behavior: Brown et al. (2020)",
        "RAG: Lewis et al. (2020)",
        "Dense retrieval + FAISS for scalable semantic search",
    ], 5)  # 6

    add_bullets_slide(prs, "End-to-End Workflow", [
        "Document collection (PDF/TXT/MD)",
        "Text extraction, cleaning, and chunking",
        "Embedding generation and FAISS indexing",
        "Top-k retrieval + context injection",
        "Local LLM answer generation",
        "Evaluation and comparison",
    ], 0)  # 7

    add_architecture_slide(prs, 1)  # 8

    add_two_column_slide(prs, "Baseline vs RAG", "Baseline LLM", [
        "Direct answer generation",
        "No external evidence grounding",
        "Lower latency",
        "Higher hallucination risk",
    ], "RAG LLM", [
        "Retrieves top-k evidence",
        "Grounded prompt construction",
        "Improved factual reliability",
        "Slight latency overhead",
    ], 2)  # 9

    add_bullets_slide(prs, "Technology Stack", [
        "Python, FastAPI, Pydantic",
        "Sentence Transformers + Transformers",
        "FAISS vector search",
        "Docker for cloud deployment",
        "Frontend: HTML/CSS/JavaScript",
    ], 3)  # 10

    add_bullets_slide(prs, "Implementation Modules", [
        "src/ingestion.py: corpus ingestion and index build",
        "src/retrieval.py: semantic search",
        "src/llm.py + src/rag_pipeline.py: response generation",
        "src/main.py: API + UI routes",
        "scripts/run_experiments.py: evaluation runner",
    ], 4)  # 11

    add_bullets_slide(prs, "API and UI Endpoints", [
        "GET /health: service health",
        "POST /ask: baseline or rag response",
        "GET /: frontend interface",
        "GET /docs: backend Swagger documentation",
        "Response includes answer, latency, and retrieved chunks",
    ], 5)  # 12

    add_bullets_slide(prs, "Experimental Strategy", [
        "Experiment 1: Baseline LLM",
        "Experiment 2: Basic RAG",
        "Experiment 3: Optimized RAG settings",
        "Run same question set across all conditions",
        "Compare quality and latency tradeoffs",
    ], 0)  # 13

    add_bullets_slide(prs, "Evaluation Metrics", [
        "Factual overlap proxy (higher is better)",
        "Hallucination proxy (lower is better)",
        "Precision@k and Recall@k",
        "Latency in milliseconds",
        "Error-case analysis by question",
    ], 1)  # 14

    add_bullets_slide(prs, "Expected Outcomes", [
        "RAG improves factual grounding over baseline",
        "Hallucination risk reduced for document-based queries",
        "Clear retrieval evidence supports trust",
        "Cloud-deployable project with reproducible artifacts",
    ], 2)  # 15

    add_bullets_slide(prs, "Six-Month Timeline", [
        "M1: Literature review and setup",
        "M2: Ingestion and indexing",
        "M3: Baseline + RAG implementation",
        "M4: Evaluation module and tuning",
        "M5: Cloud deployment and full testing",
        "M6: Final report, PPT, and viva preparation",
    ], 3)  # 16

    add_bullets_slide(prs, "Risks and Mitigation", [
        "Data quality issues -> source validation",
        "Environment/network constraints -> offline mock mode",
        "Cloud security risk -> TLS and least privilege",
        "Schedule risk -> weekly milestone tracking",
    ], 4)  # 17

    add_bullets_slide(prs, "References", [
        "Vaswani et al. (2017) - Attention Is All You Need",
        "Brown et al. (2020) - Language Models are Few-Shot Learners",
        "Lewis et al. (2020) - Retrieval-Augmented Generation",
        "Karpukhin et al. (2020) - Dense Passage Retrieval",
        "FAISS docs: https://faiss.ai/",
        "FastAPI docs: https://fastapi.tiangolo.com/",
        "Repository: https://github.com/sundar0004/-Cloud-Based-RAG.git",
    ], 5)  # 18

    add_image_slide(
        prs,
        "Frontend UI Screenshot",
        "docs/screenshots/frontend_ui.png",
        "Figure: Frontend interface for query input, answer output, and retrieved context.",
        0,
    )  # 19

    add_image_slide(
        prs,
        "Backend API Screenshot",
        "docs/screenshots/backend_api_docs.png",
        "Figure: FastAPI Swagger documentation for backend endpoint testing.",
        1,
    )  # 20

    prs.save(OUT_FILE)
    print(f"Created: {OUT_FILE}")


if __name__ == "__main__":
    build_deck()
